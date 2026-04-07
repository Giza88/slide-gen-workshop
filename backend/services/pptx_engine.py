from io import BytesIO

from pptx import Presentation


def build_presentation(slides_data: list[dict]) -> bytes:
    """
    Input: list of {"title": str, "bullets": [str, ...]}
    Output: raw .pptx bytes — no file written to disk
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content

    for slide_data in slides_data:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "Untitled")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bullet in slide_data.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
